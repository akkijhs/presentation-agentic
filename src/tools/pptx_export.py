from pptx import Presentation
from pptx.util import Inches, Pt
 
# A helper dictionary to map our AI's layout choices to PowerPoint's default layout indices
LAYOUT_MAPPING = {
    "title_slide": 0,          # Title Slide
    "title_and_content": 1,    # Title and Content
    "two_column": 3,           # Two Content
    "chart_and_text": 1        # Using 'Title and Content' as a base for charts
}
 
def export_to_pptx(presentation_data, chart_images, output_filename="AI_Presentation.pptx"):
    print(f"🎬 Starting export: {presentation_data.presentation_title}")
    # Initialize presentation (uses the default blank PowerPoint template)
    prs = Presentation()
    for slide_data in presentation_data.slides:
        # 1. Select the correct layout
        layout_idx = LAYOUT_MAPPING.get(slide_data.layout_type, 1)
        slide_layout = prs.slide_layouts[layout_idx]
        # 2. Create the slide
        slide = prs.slides.add_slide(slide_layout)
        # 3. Add the Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        # 4. Add the Body Content (Bullets)
        # Placeholders[1] is typically the main text box in standard layouts
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            # Clear default text if any, then add our bullets
            text_frame.text = "" 
            for i, bullet in enumerate(slide_data.bullets):
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0 # Bullet indentation level
        # 5. Add Speaker Notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.speaker_notes
 
        # 6. Placeholder logic for charts
        if slide_data.layout_type == "chart_and_text" and slide_data.slide_number in chart_images:
            img_path = chart_images[slide_data.slide_number]
            # Insert into PowerPoint (positioning it on the right)
            left = Inches(5.0)  
            top = Inches(2.0)
            width = Inches(4.5)
            slide.shapes.add_picture(img_path, left, top, width=width)
 
    # Save the file
    prs.save(output_filename)
    print(f"✅ Success! Presentation saved as {output_filename}")
 
# --- How to run it with the previous code ---
#
# if __name__ == "__main__":
#     # 1. Get the data from the AI Outliner (from previous script)
#     data = generate_presentation_json("AI CRM Pitch", "Retail Execs", 5)
#     
#     # 2. Pass that data to our exporter
#     export_to_pptx(data, "Pitch_Deck.pptx")